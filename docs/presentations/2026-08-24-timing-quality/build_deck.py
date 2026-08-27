"""
Janasunani 2.0 — a twelve-slide technical briefing (thirteen with the optional
Sarvam category head-to-head).

Built by CLONING slide archetypes out of the GAPG 18 August 2026 reference deck,
so type, colour, geometry and the footer lockup are that file's own shapes rather
than a reconstruction.

Archetypes reused (by reference slide number):
  S1  title
  S2  three numbered items + right-hand note + bottom line
  S3  two grouped bullet lists + bottom maroon line
  S6  three label/description rows (stripped, used as a blank shell)
  S9  three label + proportional bar + number rows
  S13 closing

Reference slides 4, 8 (portal screenshots) and 10 (native chart) are NOT used:
the screenshots carry citizen PII and the chart has no equivalent here.

Three slides are authored rather than cloned, because the reference deck contains
neither a diagram nor a table:
  slide 3  architecture.svg, embedded as true vector with a PNG fallback
  slide 4  a native table styled to the reference palette
  slide 8  the same table, on the routing-outcome diagnostics

Every figure traces to an artifact; sources live in the speaker notes.

Writing rules for slide copy: no em-dash parentheticals, no filler intensifiers
("clearly", "effectively", "genuinely"), no "not X but Y" scaffolding. One claim
per line. Say the number, then what it means.
"""

import copy
import pathlib

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.text.text import _Paragraph
from pptx.util import Inches, Pt

REF = "reference.pptx"
OUT = "Janasunani_2.0_Timing_and_Quality.pptx"
SVG = "architecture.svg"
PNG = "architecture.png"
NOTES_MD = pathlib.Path("SPEAKER_NOTES.md")

# Reference archetype indices (0-based into prs.slides)
A_TITLE, A_NUM3, A_TWOGROUP, A_ROWS3, A_BARS3, A_CLOSING = 0, 1, 2, 5, 8, 12

# Palette, read off the reference deck's own shapes.
MAROON = "7A1F2B"
DARK = "1E1F24"
BODY = "44464D"
SUBTLE = "6C6E76"
FILL = "F2F2F2"
WHITE = "FFFFFF"

# ── Sarvam head-to-head, filled in when the benchmark run lands ───────────────
# Category is per GRIEVANCE, not per page. A document spans several pages and
# carries one officer-recorded category, so `n` counts documents. Getting this
# wrong inflates n by roughly the pages-per-document ratio.
#
# Leave as None and the slide is skipped. The deck always builds, and it never
# ships a placeholder number.
#
#   n         documents graded
#   pipeline  our categoriser's agreement with the officer's recorded label
#   sarvam    Sarvam Extract's agreement with the same label
#   baseline  always guessing the single biggest category
#   title     the assertion these numbers actually support
#   closing   the line under the bars
SARVAM_HEAD_TO_HEAD = None
# Example once measured:
# SARVAM_HEAD_TO_HEAD = {
#     "n": 212, "pipeline": 0.0, "sarvam": 0.0, "baseline": 0.0,
#     "title": "", "closing": "", "notes": "",
# }

SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
SVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"


# ── cloning ───────────────────────────────────────────────────────────────────
# The notes master carries no body placeholder, so a freshly created notes slide
# has nowhere to put speaker notes. Borrow the reference deck's own notes
# placeholder and graft it onto every cloned slide.
NOTES_PLACEHOLDER = None


def clone(prs, src_idx):
    """Deep-copy every shape of a source slide onto a new slide at the end."""
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    for shp in src.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    if new.notes_slide.notes_text_frame is None and NOTES_PLACEHOLDER is not None:
        new.notes_slide.shapes._spTree.append(copy.deepcopy(NOTES_PLACEHOLDER))
    return new


def clone_shell(prs, src_idx, keep):
    """Clone an archetype and strip it back to eyebrow, title and footers.

    Used for the two authored slides. Keeping the reference's own heading and
    footer shapes means an authored slide still lines up with a cloned one to
    the pixel.
    """
    new = clone(prs, src_idx)
    for i, shp in enumerate(list(new.shapes)):
        if i not in keep:
            shp._element.getparent().remove(shp._element)
    return new


def set_text(shape, text):
    """Replace a shape's text, keeping the first run's formatting.

    Multi-line text is split on newlines into paragraphs, each cloned from the
    first paragraph so bullet/indent/spacing properties survive.
    """
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        tf.text = text
        return
    r0 = p0.runs[0]
    lines = text.split("\n")
    r0.text = lines[0]
    for extra in p0.runs[1:]:
        extra._r.getparent().remove(extra._r)
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    for line in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        para = _Paragraph(newp, tf)
        para.runs[0].text = line
        for extra in para.runs[1:]:
            extra._r.getparent().remove(extra._r)


def S(slide):
    """Shapes as a list, in document order."""
    return list(slide.shapes)


def top_anchor(shape, height_in):
    """Pin a text box to the top of a fixed band."""
    shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
    shape.height = Inches(height_in)


def textbox(slide, text, *, x, y, w, h, size, color, bold=False, italic=False,
            font="Calibri", align=PP_ALIGN.LEFT):
    """A plain text box in the reference deck's idiom."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = RGBColor.from_string(color)
    return box


# ── the SVG diagram ───────────────────────────────────────────────────────────
def add_vector_picture(slide, png_path, svg_path, *, x, y, w, h):
    """Place a picture that PowerPoint renders as vector.

    PowerPoint 2016+ stores an SVG as a normal picture whose blip carries an
    `asvg:svgBlip` extension pointing at the SVG part; the raster blip is the
    fallback for anything that cannot render SVG. python-pptx has no API for
    this, so add the PNG normally and then attach the SVG part by hand.

    The fallback matters here: this machine has neither Calibri nor Cambria, so
    the PNG is rasterised against Carlito/Caladea. PowerPoint resolves the real
    faces from the SVG and the diagram matches the rest of the deck.
    """
    pic = slide.shapes.add_picture(
        png_path, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    package = slide.part.package
    partname = PackURI(f"/ppt/media/{pathlib.Path(svg_path).name}")
    svg_part = Part(
        partname, "image/svg+xml", package, pathlib.Path(svg_path).read_bytes()
    )
    r_id = slide.part.relate_to(svg_part, RT.IMAGE)
    blip = pic._element.blipFill.find(qn("a:blip"))
    blip.append(
        parse_xml(
            f'<a:extLst {nsdecls("a", "r")}>'
            f'<a:ext uri="{SVG_EXT_URI}">'
            f'<asvg:svgBlip xmlns:asvg="{SVG_NS}" r:embed="{r_id}"/>'
            f"</a:ext></a:extLst>"
        )
    )
    return pic


def _set_cell_borders(cell, color):
    """Pin every cell edge to a hairline.

    python-pptx exposes no border API, and with the themed table style removed
    each renderer picks its own default. Setting the lines explicitly keeps
    PowerPoint and LibreOffice in agreement.
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        for existing in tc_pr.findall(qn(edge)):
            tc_pr.remove(existing)
        ln = parse_xml(
            f'<{edge} {nsdecls("a")} w="6350" cap="flat" cmpd="sng" algn="ctr">'
            f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
            f"</{edge}>"
        )
        tc_pr.append(ln)


# ── the stage table ───────────────────────────────────────────────────────────
def add_table(slide, rows, *, x, y, w, col_w, row_h, header_h):
    """A table styled to the reference palette.

    python-pptx attaches a themed table style with banding and blue accents.
    Strip it and set every fill explicitly, so the table reads as the flat
    blocks the rest of the deck uses.
    """
    total_h = header_h + row_h * (len(rows) - 1)
    gf = slide.shapes.add_table(
        len(rows), len(col_w), Inches(x), Inches(y), Inches(w), Inches(total_h)
    )
    table = gf.table
    table.first_row = False
    table.horz_banding = False

    tbl_pr = table._tbl.tblPr
    for style_id in tbl_pr.findall(qn("a:tableStyleId")):
        tbl_pr.remove(style_id)

    for i, cw in enumerate(col_w):
        table.columns[i].width = Inches(cw)
    table.rows[0].height = Inches(header_h)
    for r in range(1, len(rows)):
        table.rows[r].height = Inches(row_h)

    for r, row in enumerate(rows):
        header = r == 0
        bg = MAROON if header else (WHITE if r % 2 else FILL)
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(bg)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.10)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.name = "Calibri"
            run.font.size = Pt(13 if header else 15)
            run.font.bold = header or c == 0
            run.font.color.rgb = RGBColor.from_string(
                WHITE if header else (DARK if c == 0 else BODY)
            )
            _set_cell_borders(cell, MAROON if header else "DDDDDD")
    return table


# ── deck ──────────────────────────────────────────────────────────────────────
prs = Presentation(REF)
n_original = len(prs.slides._sldIdLst)

# Capture the reference deck's notes body placeholder before anything is removed.
for _src in prs.slides:
    _tf = _src.notes_slide.notes_text_frame
    if _tf is not None:
        NOTES_PLACEHOLDER = copy.deepcopy(_tf._txBody.getparent())
        break

# ══ 1 · Title ═════════════════════════════════════════════════════════════════
s = clone(prs, A_TITLE)
sh = S(s)
set_text(sh[0], "TECHNICAL BRIEFING")
set_text(sh[1], "Janasunani 2.0")
set_text(sh[2], "Reading, sorting and routing 1.37 million grievances")
set_text(sh[4], "Data, Policy and Innovation Centre")
set_text(sh[5], "24 August 2026")

# ══ 2 · The record ════════════════════════════════════════════════════════════
s = clone(prs, A_ROWS3)
sh = S(s)
set_text(sh[0], "THE RECORD")
set_text(sh[1], "1.37 million grievances, 6.5 million action notes")
set_text(sh[3], "1.37 million grievances")
set_text(sh[4], "What the citizen asked for, 2021 to 2025 across 30 districts. Median 19 words, mostly unique.")
set_text(sh[6], "6.5 million action notes")
set_text(sh[7], "What the officer did about it. Free text on 99.87% of rows, 1.4 million distinct values.")
set_text(sh[9], "The portal counts rows")
set_text(sh[10], "It reports how many cases are open and how many are closed. It has never opened either body of text.")
textbox(
    s,
    "Every number in this deck comes from reading one of those two.",
    x=0.7, y=6.28, w=11.9, h=0.42, size=17, color=MAROON, bold=True,
)
s.notes_slide.notes_text_frame.text = (
    "Sources: docs/ARCHITECTURE.md, docs/ROADMAP.md. Canonical counts, verified on both local "
    "SQLite and cloud Postgres, and they must match after any migration change: 1,371,288 "
    "complaints and 6,556,171 action-history rows.\n\n"
    "Caveat: the Parquet lake reads 6,548,820 action rows against the canonical 6,556,171, a "
    "0.11% shortfall tracked as issue #241. Use the canonical figure.\n\n"
    "The grievance text has a median of 19 words and 61% of it is unique. The action notes are "
    "populated on 99.87% of rows with 1,395,867 distinct normalised values. Both are opaque "
    "strings to the reporting layer.\n\n"
    "THIS SLIDE IS THE MAP. Two bodies of free text, and the deck reads them in that order. "
    "Slides 3 to 5 read the citizen's text: the stages, their speed and quality, and the "
    "department suggestion. Slides 6 to 8 read the officer's notes: how closures are recorded, "
    "and whether a different route would have closed cases faster. Say the two halves out loud "
    "here, so slide 6 lands as the turn rather than a new topic.\n\n"
    "The third row is the reason the project exists. The portal has never opened either body of "
    "text. There is also no citizen key, so every row is an island, which is why deduplication "
    "on the next slide needs the whole corpus at once."
)

# ══ 3 · Architecture ══════════════════════════════════════════════════════════
# Keep only eyebrow, title and the two footer texts from the S6 archetype.
s = clone_shell(prs, A_ROWS3, keep={0, 1, 11, 12})
sh = S(s)
set_text(sh[0], "READING THE CITIZEN'S TEXT")
set_text(sh[1], "The live path and the batch path")
add_vector_picture(s, PNG, SVG, x=0.7, y=2.08, w=11.9, h=11.9 * 450 / 1280)
textbox(
    s,
    "The live path sees one grievance. Only the batch path knows that 55,544 filings "
    "are 10,963 distinct problems.",
    x=0.7, y=6.42, w=11.9, h=0.42, size=17, color=MAROON, bold=True,
)
s.notes_slide.notes_text_frame.text = (
    "The asymmetry is the whole slide. The live path processes one grievance and can never "
    "know it duplicates another. Only the batch path compares records against each other.\n\n"
    "Duplicate matching is genuinely unavailable live, not merely unbuilt. DuplicateReview "
    "defaults to `not_indexed` with the reason 'the live submission path is not connected to a "
    "completed index'; serving/triage.py states 'duplicate matching remains slice-scoped and "
    "unavailable live'; docs/ARCHITECTURE.md records 'per-request live matching is not wired'. "
    "Do not imply a per-request duplicate check exists.\n\n"
    "Dedup figures, docs/PERFORMANCE.md section 4, Sambalpur 2024: 55,544 filings, 10,963 "
    "distinct problems, 8,560 distinct signatories, ~57 min on 2 vCPU, 16,138,623 comparison "
    "pairs. The decomposition is the argument for the index: group GOV2024999640 is 26,203 "
    "filings from ONE signatory, while DM2024854026 is 1,291 filings from 1,155. On filing "
    "counts alone those two are indistinguishable.\n\n"
    "Materialisation ~26 s at full scale, via olap/materialize.py using DuckDB.\n\n"
    "The freshness gap is by design: GET /grievance/{id} reads the transactional store and is "
    "instant; GET /history and /supervisor read the Parquet lake and lag until the next "
    "re-materialisation. Analytics never touch the transactional store.\n\n"
    "The diagram is architecture.svg in this directory, embedded as vector with a PNG "
    "fallback. Edit the SVG, re-run rsvg-convert, re-run this script."
)

# ══ 4 · Every stage ═══════════════════════════════════════════════════════════
s = clone_shell(prs, A_ROWS3, keep={0, 1, 11, 12})
sh = S(s)
set_text(sh[0], "EVERY STAGE")
set_text(sh[1], "Six stages, timed and scored")
add_table(
    s,
    [
        ["Stage", "Model", "Latency", "Quality"],
        ["OCR", "pytesseract", "5.8 s", "not measured, no ground truth"],
        ["PII redaction", "Presidio + spaCy NER", "0.06 s", "recall 0.78, n=480 spans"],
        ["Actionability triage", "TF-IDF word + char", "< 0.01 s", "recall 1.00, precision 0.81, n=57"],
        ["Categorisation", "hashing word + char SGD", "0.8 s", "top-1 47%, top-3 91%, n=3,160"],
        ["Summarisation", "BART", "6.6 s", "fact recall 66%, 8 of 26 usable"],
        ["Routing", "empirical crosswalk", "< 0.01 s", "top-1 55%, top-3 80%, n=142,181"],
    ],
    x=0.7, y=1.96, w=11.9, col_w=[2.5, 3.1, 1.5, 4.8], row_h=0.46, header_h=0.40,
)
textbox(
    s,
    "End to end: 0.13 s typed, 13.7 s scanned. OCR and summarisation are 94% of it.",
    x=0.7, y=5.28, w=11.9, h=0.34, size=15, color=DARK,
)
textbox(
    s,
    "On triage a TF-IDF model beat a frozen MuRIL probe, 13 of 13 against 9 of 13.",
    x=0.7, y=5.80, w=11.9, h=0.42, size=17, color=MAROON, bold=True,
)
textbox(
    s,
    "Latency is a mean over one scanned document, n=20. PII precision is unmeasured, though "
    "0 of 55,544 redacted records retain a shaped identifier.",
    x=0.7, y=6.20, w=11.9, h=0.26, size=12, color=SUBTLE, italic=True,
)
textbox(
    s,
    "Category classes are lopsided. A majority-class baseline scores 37% top-1 and 84% top-3, "
    "so the model adds 9 and 7 points.",
    x=0.7, y=6.43, w=11.9, h=0.26, size=12, color=SUBTLE, italic=True,
)
textbox(
    s,
    "Document measurements are one draw, Sambalpur 2024, the slice the duplicate index was "
    "built on. Categorisation, triage and routing are measured on text and records, not documents.",
    x=0.7, y=6.66, w=11.9, h=0.22, size=11, color=SUBTLE, italic=True,
)
s.notes_slide.notes_text_frame.text = (
    "MODEL OPTIONS AT EACH STAGE, if asked what else was tried.\n"
    "  OCR: pytesseract ships because it handles Odia script. DeepSeek-OCR exists in the repo "
    "but is English-only in practice and GPU-bound. Sarvam Vision is the third option and is "
    "slide 11.\n"
    "  Actionability: a frozen MuRIL probe was the alternative and lost, 9 of 13 review cases "
    "against TF-IDF's 13 of 13, at 86.0% accuracy against 94.7%. The cheap model won, and that "
    "is the procurement point.\n"
    "  Categorisation: a hashing word and character SGD ships. The DSI legacy MuRIL figure of "
    "71.04% is NOT a comparison; it was measured on typed subject lines with a different split "
    "and issue #127 warns against putting the two side by side.\n"
    "  Routing: the empirical crosswalk ships. An opt-in empirical-Bayes incidence provider "
    "exists behind JANASUNANI_ROUTER=incidence and falls back safely when it cannot answer.\n"
    "  Summarisation: local BART, bart-large-cnn rev 37f520fa. No alternative has been "
    "benchmarked.\n\n"
    "SPEEDS. outputs/benchmark/latency.json, run 2026-08-10T23:14:58Z at git sha 24ab193, "
    "is_fake_timing false. Document path, n=20 over 10 clusters. Means: OCR 5.833, summarise "
    "6.550, categorise 0.778, redact 0.055, detect PII 0.021, route 0.00048, triage 0.00026. "
    "Summarise plus OCR is 12.383 s of the 13.244 s mean run, 93.5%. End to end: typed p50 "
    "0.133 s (n=40), PDF p50 13.661 s (n=20). Live API warm POST median 4.44 s (n=8), cold "
    "start 19.4 s. Measured on an arm64 laptop, not the deployment box.\n\n"
    "Four stages (format classifier, page type, pii, spam) were never separately instrumented "
    "and carry n=0. They are omitted rather than shown as zero.\n\n"
    "PROVENANCE, not yet reflected in the numbers above. A 224-document real-document sample "
    "(451 pages, 30 categories, nesting the earlier 92-ticket Sarvam sample by ticket id) is "
    "staged for Sambalpur 2024, the slice the duplicate index was built on. Re-running the "
    "latency harness (feat/benchmark-real-documents, load_staged_documents) against it was "
    "deliberately deferred to avoid CPU contention with a concurrent Sarvam evaluation; the n=20 "
    "synthetic-fixture numbers above are unchanged pending that run.\n\n"
    "QUALITY, with intervals.\n"
    "Redaction: overlap recall 0.779, coverage 0.783, exact 0.550 on 480 hand-marked gold "
    "spans across 89 pages and 50 documents. Per entity: Aadhaar 0.857 (n=7), phone 0.828 "
    "(n=29), name 0.777 (n=404), email 0.750 (n=40). Corpus scan: 0 of 55,544.\n"
    "Triage: n=57 held-out, accuracy 94.74% (85.63-98.19), review recall 13/13 = 100% "
    "(77.19-100), false-flag rate 3/44 = 6.82% (2.35-18.23). TF-IDF word+char beat a frozen "
    "MuRIL probe 13/13 against 9/13.\n"
    "Category: top-3 90.89%, top-1 46.55%, n=3,160, chronological 2024 split, "
    "exact-text-group-disjoint, macro-F1 36.5%, ECE 26.4%.\n"
    "  AGAINST A TRIVIAL BASELINE, computed from the per-class supports in "
    "outputs/evaluation/categorization_historical_v1.json. The test set is lopsided: Social "
    "Welfare 1,179, Housing 743, Miscellaneous 724, together 2,646 of 3,160. Always guessing "
    "the single biggest scores 37.3% against the model's 46.6%, a lift of 9.2 points. Always "
    "guessing the biggest three scores 83.7% against the model's 90.9%, a lift of 7.2 points. "
    "Say this if anyone treats 91% as the headline.\n"
    "  Per class the spread is wide. Best: Land Matters F1 62.4%, Energy 62.3%, Social Welfare "
    "62.1%. Worst: General 4.6%, Public Utility 11.1%, Financial Assistance 12.1%. "
    "Miscellaneous has 724 cases and 11.7% recall, because a catch-all class has no signature "
    "to learn. No class sits at F1 zero.\n"
    "  The DSI reference of 71.04% for MuRIL is NOT a comparison. It was measured on typed "
    "subject lines with a different split and issue #127 warns against putting the two side by "
    "side. We have not re-run that model on this split, so no head-to-head exists.\n"
    "Summary: critical-fact recall 65.48% over 84 facts, 8 of 26 usable unedited, residual PII "
    "in output 4/26 = 15.38%. One judge, not an officer. All four coherent Odia cases were "
    "skipped by an English-only gate.\n"
    "Department: top-3 79.68% / top-1 54.96% on informative categories (n=142,181); 69.04% / "
    "45.14% across all eligible (n=208,267). Untouched 2025 test year. This measures agreement "
    "with where cases were historically sent, not jurisdictional correctness.\n\n"
    "WITHDRAWN, do not use: the in-sample crosswalk figures 60.9 / 67.5 / 72.8; PII coverage "
    "49.6% and name 0.44; MuRIL 71.04% as a current number.\n\n"
    "If asked about routing time savings: an estimated 11 to 23 day gain held on validation "
    "2024 and failed on the untouched 2025 test year at -2.35 days against a standard error of "
    "3.50. It was withdrawn on 23 August, commits 879c24c and 365e3b4, and four artifacts are "
    "archived as do-not-cite. The direct and doubly-robust estimators disagree by 33 days on "
    "the test year, which is the diagnosis: no overlap to estimate on."
)

# ══ 5 · Routing ═══════════════════════════════════════════════════════════════
s = clone(prs, A_ROWS3)
sh = S(s)
set_text(sh[0], "SUGGESTING A DEPARTMENT")
set_text(sh[1], "The department suggestion: an empirical crosswalk")
set_text(sh[3], "Learned from history")
set_text(sh[4], "5,084 keys built from where past grievances were actually sent. Right department in its top three 80% of the time on a year we held back.")
set_text(sh[6], "Where it falls back")
set_text(sh[7], "No master table links a category to a department, so the keys had to be learned. When no key matches it drops to matching on name.")
for _dead in (sh[8], sh[9], sh[10]):
    _dead._element.getparent().remove(_dead._element)
textbox(
    s,
    "All of it learns where grievances were sent. None of it knows where one should go.",
    x=0.7, y=5.05, w=11.9, h=0.42, size=17, color=MAROON, bold=True,
)
s.notes_slide.notes_text_frame.text = (
    "The ladder is janasunani/routing/provider.py. Three modes: ROUTER_DEFAULT 'crosswalk' is "
    "the shipped path and runs crosswalk, then mapping tables, then generic fallback; 'rules' "
    "skips the crosswalk and reproduces pre-#33 behaviour, useful to isolate the crosswalk in "
    "a comparison; ROUTER_INCIDENCE 'incidence' serves a checksummed empirical-Bayes artifact "
    "with the same ladder underneath it.\n\n"
    "RUNG 1, the crosswalk. Artifact janasunani/routing/reference/routing_crosswalk.json: 34 "
    "by-category keys, 257 by-subcategory, 971 by-category-district, 5,084 on the full key. "
    "Held-out performance, outputs/evaluation/routing_historical_{informative,all}.json: top-3 "
    "79.68% and top-1 54.96% on informative categories (n=142,181); 69.04% and 45.14% across "
    "all eligible (n=208,267). Train 2021-23, validate 2024, final refit on train plus "
    "validation, test on an untouched 2025. Selected alpha=100 with a one-year history window.\n\n"
    "Below the crosswalk sit the ORTPSA master tables and a generic fallback. They are off "
    "this slide because they rarely answer: intCategoryGrp is NULL on all 62 categories, so no "
    "category-to-department link exists and MappingRouter can only bridge by exact name. That "
    "absence is why the crosswalk had to be learned from history rather than read off a table. "
    "Mention it only if someone asks what happens when the crosswalk has no key.\n\n"
    "THE OUTCOME MODEL IS NOT THIS SLIDE. Everything here predicts where a grievance was sent. "
    "Whether sending it elsewhere would have resolved it better is slides 7 and 8, it is "
    "research-only, and no serving provider reads it. Keep the two apart when answering.\n\n"
    "RUNG 3, the trained model. Opt-in via JANASUNANI_ROUTER=incidence, artifact checksummed, "
    "and IncidenceRoutingProvider falls through to the crosswalk and rules when a lookup "
    "fails, logging rather than failing silently.\n\n"
    "THE CAVEAT THAT MATTERS, and it applies to all three rungs. Every one of them measures "
    "agreement with the historical destination, not jurisdictional correctness. A "
    "correct-authority adjudication does not exist. Roughly 300 closed cases read by hand "
    "would settle it.\n\n"
    "Macro-F1 is weak, 25.2% informative and 19.8% all eligible, and about a dozen departments "
    "sit at F1 zero. The top-three framing is doing real work here.\n\n"
    "WITHDRAWN, do not use: the in-sample crosswalk figures 60.9 / 67.5 / 72.8, which are "
    "resubstitution. Any routing time saving is withdrawn too; slide 8 carries that warning."
)

# ══ 6 · Closing a case costs nothing ══════════════════════════════════════════
s = clone(prs, A_BARS3)
sh = S(s)
set_text(sh[0], "READING THE OFFICER'S NOTES")
set_text(sh[1], "Closures on the six-template disposal ladder")
set_text(sh[2], "How 776,922 closures were recorded, on the six standard remarks.")
ladder = 776_922
rungs = [
    (sh[3], sh[4], sh[5], "Disposed, no action claimed", 472_782),
    (sh[6], sh[7], sh[8], "Disposed with action taken", 280_887),
    (sh[9], sh[10], sh[11], "Beneficiary benefited", 23_253),
]
for label_sh, bar_sh, num_sh, label, count in rungs:
    set_text(label_sh, label)
    set_text(num_sh, f"{count:,}")
    width = max(1.35, 8.0 * count / 472_782)
    bar_sh.width = Inches(width)
    num_sh.width = Inches(width - 0.3)
set_text(sh[15], "Median days to close: 46 claiming no action, 54 claiming action taken")
set_text(sh[12], "The fastest way to close a case is to do nothing. Speed alone cannot be the target.")
s.notes_slide.notes_text_frame.text = (
    "THIS IS THE TURN slide 2 announced. Everything up to here read the citizen's text. From "
    "here the source is the officer's action notes, the second of the two bodies. Say so; it "
    "is one sentence and it stops this slide reading as a new topic.\n\n"
    "It sets up the next slide and it is also the strongest finding in the record.\n\n"
    "NUMBERS. outputs/findings/closure_finding.md. Of 776,922 complaints closed on one of six "
    "governed disposal templates, 472,782 used the rung that records no action. That is 60.9% "
    "of ladder closures and 39.1% of all 1,209,144 resolved complaints. State the denominator, "
    "because it moves the number by half. The ladder covers 64.25% of resolved cases; the rest "
    "carry non-standard text.\n\n"
    "WHY IT MATTERS FOR ROUTING. From the design document, section 2: 'The first instinct is to "
    "minimise time to closure. This fails, instructively. Closure is an action the officer "
    "controls directly, and the optimal behaviour under that objective is to close every case "
    "immediately without doing anything. This is not hypothetical: the most common closing "
    "remark in the system records disposal with no action claimed.' That is why the model "
    "minimises duration SUBJECT TO a correctness constraint rather than minimising duration.\n\n"
    "It is also the weak point. The constraint is derived from the closing remark after the "
    "fact, so the population is conditioned on resolution and does not identify intake-time "
    "actionability. It is on the unresolved list at slide 8, and one reason the estimate there "
    "did not hold.\n\n"
    "THE CAVEAT THAT MUST TRAVEL. This is descriptive and is not a failure rate. A correct "
    "closure and a premature one look identical in this record. Do not present it as an office "
    "league table; report it at state level. The bare share RISES with work done, 58.4% at "
    "three to five action steps and 64.8% at six or more, so whatever drives the choice of "
    "phrase, it is not that nothing happened.\n\n"
    "The benefit flag runs backwards to the ladder: 9.2% of no-action closures carry it against "
    "6.2% of action-taken closures. Never report it as satisfaction.\n\n"
    "If asked for a bounded review set: 8,974 complaints were created and closed within two "
    "days on a bare disposal, 1.9% of bare disposals."
)

# ══ 7 · What was built ════════════════════════════════════════════════════════
# The archetype's own maroon 1/2/3 numerals carry the sequence, so the three
# stages need no authored shapes. Its bottom line at sh[16] is dropped: the
# deck's idiom puts the maroon conclusion last, and the live/research
# distinction has to sit above it.
s = clone(prs, A_NUM3)
sh = S(s)
set_text(sh[0], "WOULD ANOTHER ROUTE WORK BETTER?")
set_text(sh[1], "How the routing evaluation was built")
set_text(sh[3], "Estimate time to disposal on every route")
set_text(sh[5], "Ridge and gradient boosting, cross-fitted, capped at 365 days, under each historically supported department and role chain.")
set_text(sh[8], "Safeguard the action rate")
set_text(sh[10], "A calibrated classifier for whether substantive action is taken. A faster route is admissible only where that rate holds.")
set_text(sh[13], "Verify against outcomes we observed")
set_text(sh[15], "Direct outcome regression against augmented IPCW, Hajek-normalised, with censoring weights and a district-by-year cluster bootstrap.")
sh[16]._element.getparent().remove(sh[16]._element)
textbox(
    s,
    "LIVE: predicts where similar grievances were historically sent.",
    x=0.7, y=5.90, w=11.9, h=0.26, size=13, color=SUBTLE,
)
textbox(
    s,
    "RESEARCH ONLY: asks where they would have resolved better.",
    x=0.7, y=6.16, w=11.9, h=0.26, size=13, color=SUBTLE,
)
textbox(
    s,
    "A complete offline evaluation. It was never connected to live routing.",
    x=0.7, y=6.52, w=11.9, h=0.42, size=17, color=MAROON, bold=True,
)
s.notes_slide.notes_text_frame.text = (
    "This slide is the design. The next one is why it is not live. Do not merge them.\n\n"
    "TREATMENT is the department and the complete role chain chosen jointly at assignment, the "
    "intention to route rather than the final state. ESTIMAND is capped RMST at 365 days.\n\n"
    "THE CONSTRAINT IS THE POINT, and it is what slide 6 was for. The objective minimises "
    "expected capped duration SUBJECT TO the action rate holding at its historical level. It "
    "is never conditional on the realised outcome, which would condition on a post-treatment "
    "variable. Unconstrained duration minimisation selects for closing cases untouched, "
    "because closure is an action the officer controls directly. Design document, section 2: "
    "'The first instinct is to minimise time to closure. This fails, instructively.'\n\n"
    "THE THREE STAGES. Ridge and gradient boosting are carried separately rather than "
    "ensembled, so model-family disagreement stays visible. It later mattered. The action "
    "classifier is isotonic-calibrated, test ECE 0.0615 and AUC 0.7425. The threshold is the "
    "smallest floor meeting the constraint. Verification puts the direct outcome regression "
    "against the augmented estimator on the arm where history happened to match the proposed "
    "route, which is the only place a model prediction can be checked against an observed "
    "outcome.\n\n"
    "THE FIREWALL, if asked what is running today. Historical-route prediction is live and is "
    "slide 5. This work has no serving provider, no egress route, and no deployment "
    "recommendation. The two answer different questions: where similar grievances were sent, "
    "against where they would have resolved better.\n\n"
    "SOURCES. Design of record is docs/experiments/routing-outcome-model.tex, 42 pages. "
    "Appendices A to F carry the efficient influence function, Neyman orthogonality, the "
    "censoring argument and a marginal sensitivity model. He may ask whether they exist. They "
    "do. Plain-prose companion is janasunani/experiments/routing_outcome/README.md. Live path "
    "is janasunani/routing/provider.py."
)

# ══ 8 · Why it is not live ════════════════════════════════════════════════════
s = clone_shell(prs, A_ROWS3, keep={0, 1, 11, 12})
sh = S(s)
set_text(sh[0], "WHY IT IS NOT LIVE")
set_text(sh[1], "Validation and untouched-year diagnostics")
textbox(
    s,
    "The same pipeline, refit and rerun on a year it had never seen.",
    x=0.7, y=2.02, w=11.9, h=0.36, size=17, color=BODY,
)
add_table(
    s,
    [
        ["What we measured", "2024, held out", "2025, untouched"],
        ["Direct estimate, outcome regression", "+24.5 days", "+30.5 days"],
        ["Cross-fitted AIPW", "+26.8 days (SE 4.0)", "\u22122.3 days (SE 3.5)"],
        ["Falsification: predicting the arm actually observed", "5% error", "118% error"],
        ["Proposed route in observed support", "38% of cases", "16% of cases"],
        ["Effective sample as a share of n", "0.19", "0.07"],
    ],
    x=0.7, y=2.55, w=11.9, col_w=[6.0, 2.95, 2.95], row_h=0.46, header_h=0.40,
)
textbox(
    s,
    "DECISION   Withdraw the saving estimate, the correctness threshold and the routing recommendation. Keep the research offline.",
    x=0.7, y=5.42, w=11.9, h=0.34, size=15, color=DARK,
)
textbox(
    s,
    "Censoring at the cut runs 9% then 34%. The 2025 cohort has seven months against a 365-day outcome, so non-replication and immaturity are not separable here.",
    x=0.7, y=5.92, w=11.9, h=0.26, size=12, color=SUBTLE, italic=True,
)
textbox(
    s,
    "Unresolved: assignment-field provenance, intake-time actionability, interference since routing changes the queue, mature follow-up, a governed pilot.",
    x=0.7, y=6.16, w=11.9, h=0.26, size=12, color=SUBTLE,
)
textbox(
    s,
    "We withdrew the estimate, not the research question.",
    x=0.7, y=6.52, w=11.9, h=0.42, size=17, color=MAROON, bold=True,
)
s.notes_slide.notes_text_frame.text = (
    "Ninety seconds. The table is the argument; do not narrate every row.\n\n"
    "THE HEADLINE. The two estimators diverge by 33 days on 2025, 30.53 against -2.35, against "
    "2.3 days on 2024. Under correct specification they estimate the same functional. "
    "Divergence of that size disqualifies the pair, it does not offer a menu to choose "
    "from.\n\n"
    "DO NOT ATTRIBUTE THE FAILURE TO POSITIVITY ALONE, and do not let the title be heard "
    "that way. Weak overlap is real: support falls to 16% and effective sample to 7% of n. "
    "But row 3 is measured on the arm each case actually received, so it extrapolates "
    "nowhere, and a 118% error there is an outcome-model or cohort-drift problem rather "
    "than an overlap one. Several threats fail at once and this design cannot separate "
    "them. The honest claim is that every check refused the result, not that one of them "
    "explains it.\n\n"
    "ROW 3 IS THE ONE TO DWELL ON. It is the falsification that localises the failure. On the "
    "historical arm, the one arm where truth is observed, the direct estimator predicts 108.69 "
    "against an observed IPCW mean of 49.74. It is 118% wrong where it can be checked, and it "
    "was 5% wrong on 2024. AIPW down-weights the cells driving that error and lands at an "
    "effective sample of 8,291 on n=113,535.\n\n"
    "SOURCES, live 19 August refit, both under outputs/experiments/routing_outcome/. "
    "ope_val_ridge.json (n=450,567) and ope_test_ridge.json (n=113,535). Row 1 delta_direct, "
    "24.50 and 30.53. Row 2 delta_dr, 26.77 (SE 4.04) and -2.35 (SE 3.50). Row 3 v_direct "
    "against v_dr on the historical arm, 97.84 against 93.22 then 108.69 against 49.74. Row 4 "
    "overlap.match_rate, 0.380 and 0.156. Row 5 ess over n, 87,569 and 8,291 against the "
    "respective sample sizes. Censoring is censoring_rate_of_split, 0.092 and 0.344. Governed "
    "aggregate is docs/experiments/routing-outcome-evidence-2026-08-19.json, model commit "
    "05f50da, a DVC dependency of the benchmark bundle. Status line in "
    "docs/QUALITY_BENCHMARKS.md, section 'Routing, outcome not destination': no routing gain "
    "established.\n\n"
    "'Which estimator do you believe?' Neither, on the test year. Row 3 disqualifies the "
    "direct one. The augmented one is honest about how little support remains and cannot "
    "resolve anything at that effective sample. The test year cannot answer the question.\n\n"
    "'Is this non-replication or immature follow-up?' Not separable in this design. The "
    "snapshot is 2025-07-30, so the cohort has at most seven months against a 365-day outcome "
    "and censoring runs 3.1%, 9.2%, 34.4%. It needs a later snapshot, not a better "
    "estimator.\n\n"
    "'A different outcome model?' Boosting gives 12.40 on validation against ridge's 26.77, "
    "and 0.15 on test against -2.35. The families disagree by a factor of two on the year that "
    "looked positive and agree on zero for the year that did not. The 2024 gain was never "
    "stable.\n\n"
    "ON THE UNRESOLVED LINE. Unconfoundedness is invoked on an information set smaller than "
    "the assigning officer's; congestion and trailing destination performance are absent from "
    "the covariates. Treatment provenance is open, since dept_id and vchAllEscUser may record "
    "final state rather than initial assignment. No interference is assumed, and the design "
    "document says that assumption 'in a queueing system is false'. Route everyone to the fast "
    "office and it stops being fast. That bounds any future version of this exercise.\n\n"
    "THE CORRECTNESS FRONTIER was withdrawn separately, issue #284, pending corrected "
    "regeneration. Its augmented score was normalised over all evaluation rows before being "
    "restricted to rows carrying an observed correctness label, while the direct value used "
    "the full population. tau_star is null.\n\n"
    "WITHDRAWN, do not use: any routing time saving, including the 11 to 23 day band. Do not "
    "present it as conservative, provisional or pending. There is no estimate.\n\n"
    "DO NOT USE routing-outcome-ope-val-ridge-2026-08-13.json. Its validation match rate of "
    "0.536 and its +20.1 belong to the withdrawn run and it sits in docs/experiments/"
    "superseded/. Quoting it to dramatise the collapse re-cites the thing we retracted.\n\n"
    "What we are NOT saying: that routing gains do not exist. Say no routing gain established. "
    "Never say routing does not work. This design on this data cannot detect one."
)

# ══ 9 · What the record can now be asked ════════════════════════════════════
s = clone(prs, A_NUM3)
sh = S(s)
set_text(sh[0], "WHAT THE RECORD CAN NOW BE ASKED")
set_text(sh[1], "Duplicate-adjusted workload, spikes and themes")
set_text(sh[3], "Duplicate-adjusted workload.")
set_text(sh[5], "55,544 filings in the slice are 10,963 distinct problems from 8,560 citizens. The portal reports the first number.")
set_text(sh[8], "A spike with the cause attached.")
set_text(sh[10], "Every spike carries filings, distinct problems and distinct citizens. A campaign and two hundred separate problems count the same.")
set_text(sh[13], "Local issue themes.")
set_text(sh[15], "Grouped by what complaints are about, not the category assigned at intake. Concentrated in one place and rising is the alert worth acting on.")
set_text(sh[16], "If an analyst with database access could produce a number in a day, we do not present it as something the system enables.")
s.notes_slide.notes_text_frame.text = (
    "THE HOUSE RULE, janasunani/analytics/README.md. An insight is something the record already "
    "contained and nobody had queried. A capability is something no existing dashboard could "
    "produce. Presenting the first as the second is the failure mode, ROADMAP section 5.3. All "
    "three items on this slide are capabilities: none is a query an analyst with database access "
    "could run against the existing dashboards in a day.\n\n"
    "ITEM 1, workload. 55,544 filings in Sambalpur 2024 are 10,963 distinct problems from 8,560 "
    "citizens. docs/PERFORMANCE.md section 4, production run on the deployment CPU box: ~57 "
    "minutes on 2 vCPU, 16,138,623 comparison pairs. Re-run attempted locally on this branch "
    "(janasunani-publish-workload): the only dedup index reachable outside the deployment box is "
    "a 20-ticket local smoke-test population, which verifies the digest-guard plumbing end to end "
    "(20 filings, 14 distinct, source_snapshot_id checked) but is not the governed slice, and its "
    "numbers are not quoted here.\n\n"
    "ITEM 2, spike. A campaign is not a false spike. Spike detection must not run on de-duplicated "
    "counts, and every spike is reported as three numbers, filings / distinct problems / distinct "
    "citizens, never collapsed to one. outputs/findings/spike.md still reads NULL, status 'pending "
    "dedup index': that is a hardcoded literal in analytics/sql/spike.sql (the spike_decomposition "
    "view), not a live check, so it reads the same whether or not a dedup index exists. The "
    "dedup-aware path that does compute real numbers (compute_spike_with_dedup, "
    "--publish-aggregates) runs, but is not scoped to a district/year before it picks its top "
    "candidate: on this branch it decomposed a Bargarh 2023 spike against a Sambalpur "
    "2024-scoped dedup index and silently returned no reduction at all (filings = distinct "
    "problems = distinct citizens). Filed as tech debt, not fixed on this branch.\n\n"
    "ITEM 3, themes. janasunani-publish-themes crashed on this slice, ColumnNotFoundError "
    "'filings', because render_markdown assumed the full theme schema and the 'too few rows in "
    "this category' early-exit path returns a different one. Fixed on this branch "
    "(tests/test_themes.py adds coverage for both early-exit shapes). It now reports 'No themes "
    "computed (insufficient data for themes)' for Housing, Sambalpur 2024, honestly, given only "
    "20 locally-available redacted records rather than crashing.\n\n"
    "PROVENANCE DISCIPLINE. janasunani/serving/intelligence.py reports each dashboard panel as "
    "Recorded or Unavailable with a stated reason, never as a silent zero (tests/"
    "test_intelligence.py). Marts are portable SQL handed to the department to run for "
    "themselves (janasunani/analytics/marts.py). Findings emit aggregates only, never a row of "
    "citizen writing (janasunani/analytics/README.md house rules).\n\n"
    "FOR CONTEXT IF ASKED. Closure ladder: 472,782 of 776,922 disposed with no action claimed "
    "(docs/FINDINGS.md), against 1,209,144 total resolved complaints. Confirmed duplicates, the "
    "manual-process baseline: 37,299 officer-confirmed action rows, 21,117 already-taken-up plus "
    "16,182 duplicate-copy (docs/PERFORMANCE.md section 4). The MinHash increment is reported "
    "separately and is never merged into that baseline. An earlier mart run mis-quoted 18,432 "
    "duplicates from a template-matching defect; docs/PERFORMANCE.md says explicitly not to "
    "quote it, and this deck does not."
)

# ══ 10 · The limits ═════════════════════════════════════════════════════════
s = clone(prs, A_NUM3)
sh = S(s)
set_text(sh[0], "THE LIMITS")
set_text(sh[1], "What we cannot measure")
set_text(sh[3], "We cannot say how accurately it reads a page.")
set_text(sh[5], "Nobody has hand-typed a set of pages to check against. No system can be scored, ours or anyone else's.")
set_text(sh[8], "We cannot say how often redaction hides too much.")
set_text(sh[10], "We marked what should be hidden. We never marked what should not, so a miss and an over-redaction look the same.")
set_text(sh[13], "We cannot say it saves officer time.")
set_text(sh[15], "The 10 to 15 minute baseline is self-reported and was never timed. No trial has been run.")
set_text(sh[16], "Every figure in this deck was produced by us, on our own data. None of it has been checked by an officer.")
s.notes_slide.notes_text_frame.text = (
    "Do not cut this slide for time. It is the one that makes the rest credible.\n\n"
    "1. OCR ground truth was never commissioned and has no owner, issue #53. It cannot be "
    "produced by an agent, because it is the answer key an agent would be scored against.\n\n"
    "2. 824 predicted spans against 480 marked. The recogniser over-fires, and the gold set "
    "cannot separate a missed label from an over-redaction. Filed as open.\n\n"
    "3. The officer-hours figure that circulates, 201,000 to 302,000 across 1,209,144 resolved "
    "cases, is a denominator and not a saving. It is the size of the prize, not anything "
    "realised.\n\n"
    "Two more we do not claim, if asked: no gain from duplicate detection beyond the 37,299 "
    "repeats officers already confirmed, and no OCR accuracy result for the outside option on "
    "slide 11. Its latency is measured and its category was graded; it is transcription "
    "accuracy specifically that has no referee, because divergence never says which system is "
    "right and no hand-transcribed answer key exists (#53).\n\n"
    "If someone asks why so little is claimed: because the alternative is claiming things we "
    "cannot defend, and this deck has to survive the room checking it."
)

# ══ 11 · Sarvam ═════════════════════════════════════════════════════════════
s = clone(prs, A_TWOGROUP)
sh = S(s)
set_text(sh[0], "THE OUTSIDE OPTION: SARVAM")
set_text(sh[1], "Sarvam: fields returned, latency and divergence")
set_text(sh[2], "What it does")
set_text(sh[4], "Two calls a page, ₹0.50 to read it and ₹1.00 to pull fields out: text, category, summary.")
set_text(sh[6], "It reads 22 Indian languages including Odia. Ours reads English.")
set_text(sh[12], "It reads handwriting. Ours does not.")
set_text(sh[7], "What we measured, on 200 pages of Sambalpur 2024")
set_text(sh[14], "6 s to read a page, 11 s to extract the fields, at ten requests a minute.")
set_text(sh[16], "The two systems differed on 197 of 198 pages. Sarvam returns 1.6 times the characters.")
set_text(sh[8], "Different is not better, and using it sends citizen documents outside our control.")
textbox(
    s,
    "Capped at ten requests a minute, and two calls a page: this run took 166 minutes. "
    "At the cap the English corpus is thirteen days; at the rate this run actually achieved, about eight weeks. "
    "It returned a category on 11 of 87 grievances and matched none, because our v1 schema never listed the categories.",
    x=0.7, y=6.62, w=11.9, h=0.36, size=11, color=SUBTLE, italic=True,
)
s.notes_slide.notes_text_frame.text = (
    "The provider is Sarvam. Two endpoints, billed separately: digitise at ₹0.50 a page "
    "returns text and layout; extract at ₹1.00 a page returns schema-driven fields. Both is "
    "₹1.50. Source: janasunani/evaluation/pricing.py, checked against the Sarvam dashboard "
    "2026-08-07. Our local pipeline is ₹0.00 a page.\n\n"
    "THE EXTRACT FIELDS, which is the 'does more' claim: grievance_category, summary and "
    "grievance_text. That is our OCR, categoriser and summariser in one call. "
    "janasunani/evaluation/sarvam_grievance_schema.py, schema v2, pinned so a later edit "
    "cannot silently move a headline number. docs/DELIVERY.md:163: 'our pipeline has no "
    "equivalent'. The 2026-08-25 run below was made against v1, which had a fourth field, "
    "district, and no enum on the category — both of which is why it graded the way it did. "
    "v2 is what a corrected run would use.\n\n"
    "LANGUAGE. Sarvam Vision lists all 22 scheduled languages plus English, Odia among them. "
    "There is also a separate transliteration API for romanized Odia (od-IN), which we do not "
    "solve at all today. Our own summariser skipped all four coherent Odia cases through an "
    "English-only gate, and non-English text is downgraded to Uncategorized.\n\n"
    "WHAT WE RAN, 2026-08-25. A 200-page stratified draw from Sambalpur/2024, both arms, "
    "₹300 at list price, 87 grievances across 30 categories. 200 pages processed, 198 scored, "
    "2 SarvamError failures. outputs/sarvam_run_2026-08-25/. Two earlier runs are superseded: "
    "a 5-page validation and a 300-page run that died at 65 pages on credit exhaustion.\n\n"
    "LATENCY, and this is the number the deck lacked until now. Derived from submission to "
    "result in the egress audit log, outputs/sarvam_run_2026-08-25/sarvam_audit.sqlite, "
    "excluding the later recovery re-reads. Digitise: median 6 s a page, mean 9 s, p90 16 s, "
    "n=200. Extract: median 11 s, mean 13 s, p90 19 s, n=198. Polling is 5 s quantised, so "
    "treat these as 5 s resolution. DO NOT COMPARE THIS TO OUR 5.8 s. That figure is from an "
    "n=20 synthetic-fixture run, not from documents, so the two are not the same class of "
    "input and a per-page comparison between them means nothing. A real-document local "
    "measurement is #318; until it lands there is no like-for-like latency comparison to make. "
    "THROUGHPUT is the real constraint, and note it is two "
    "calls a page, not one: --arm both submits digitise and extract separately "
    "(sarvam_evaluate.py:596 and :604), billed ₹0.50 and ₹1.00. So the 10-requests-a-minute "
    "cap, which does not rise with the plan tier, is a ceiling of 5 pages a minute for both "
    "arms. The run averaged 1.2 pages a minute end to end, about a quarter of that ceiling. "
    "Quote both numbers, because the gap between them is the point. The 96,469-page English "
    "corpus is 192,938 requests on both arms: 13 days at the cap (192,938 / 10 a minute), "
    "which is a floor and assumes perfect saturation, and about 56 days at the 1.2 pages a "
    "minute we measured (96,469 / 1.2). The calls are sequential and blocking, so the "
    "measured rate is the one to plan against. Digitise alone is about seven days at the cap. "
    "If asked why the gap: nothing here is batched or concurrent, and that is a property of "
    "this runner, not of the provider.\n\n"
    "THE CATEGORY GRADE, and say plainly that it measures our schema and not the provider. "
    "Sarvam returned a category on 11 of 87 grievances and matched the officer's label on 0 of "
    "those 11. What came back was a subject line rather than a taxonomy label: the specific "
    "relief the petitioner asked for, or the name of an administrative wing. The v1 schema "
    "described the field as the complaints taxonomy and then never listed it, so the model had "
    "nothing to choose from. Schema v2 adds the 35-label enum and is now the default. A "
    "corrected run is extract-only, about ₹200.\n\n"
    "WHY DISTRICT IS GONE FROM v2, if asked. It was the best-populated field in the run, 100 of "
    "198 pages, because it sits on the letterhead. It is also a structured column on complaints, "
    "recorded at intake and known for every grievance, so the extraction budget went to the one "
    "field we already had while the field we needed came back on 11. Reinstate it only as an "
    "explicit intake-accuracy cross-check.\n\n"
    "TRANSCRIPTION. Normalised exact-text divergence 0.995, 95% CI [0.985, 1.005], n=198 pages "
    "clustered on 87 tickets: the two systems differed on 197 of 198. Sarvam returns 295,117 "
    "normalised characters against pytesseract's 189,497, a ratio of 1.56. Divergence says they "
    "disagree, never who is right, and no hand-transcribed answer key exists (#53).\n\n"
    "SUMMARY WAS NOT GRADED, and could not have been. It was only ever scoped as divergence "
    "against BART with no gold referee, and summary_divergence is the same function as "
    "divergence_rate under another name (sarvam_scorecard.py:234). Even fully run it could "
    "only have said the two summaries differ, never which was better.\n\n"
    "  The superseded 300-page run returned 61 extract jobs that were compared against "
    "nothing, because that sample was never joined to the recorded category. The join is what "
    "--join-metadata does, and the 2026-08-25 run used it — which is why there is a category "
    "grade to report at all this time.\n\n"
    "  A schema bug returned HTTP 400 on every extract submission until 2026-08-09, which is "
    "why the 5-page validation run has no extract output at all. Every test mocked the "
    "transport, so no test could see the 400.\n\n"
    "NOT MEASURED, do not claim: OCR accuracy, summary quality, actual billed cost (every "
    "rupee figure is list price), observed language split, handwritten versus printed split. "
    "Latency and category accuracy ARE now measured — see above, and carry the caveat with "
    "the category number rather than quoting it bare.\n\n"
    "GOVERNANCE. Trust tier authorized-external. Authorisation is a GoO-Sarvam MoU with "
    "sign-off from the Additional Chief Secretary, Electronics & IT, accepted 2026-08-07, on "
    "the basis that no state statute currently governs the transfer. All three provider "
    "controls remain UNVERIFIED: retention terms, encryption in transit, encryption at rest. "
    "Authorisation and verification are recorded separately on purpose. One module may make "
    "the call, janasunani/egress/, every attempt is audit-logged, and a kill switch falls back "
    "to local pytesseract.\n\n"
    "COST AT SCALE, projected list price: ₹48,000 to digitise the 96,469-page English corpus, "
    "₹145,000 for both endpoints, ₹8,050 to push 1.37M subjects through the 105B text model. "
    "At 10 requests a minute, which does not rise with the plan tier, and two calls a page on "
    "both endpoints, the full corpus is thirteen days of continuous calling at the cap and "
    "about eight weeks at the throughput actually observed; digitise alone is about seven days "
    "at the cap. It is a measurement instrument, not a backfill path. Do not quote "
    "₹700; that was priced on the withdrawn 30B model."
)

# ══ 12 · Category head to head (only when measured) ═════════════════════════
if SARVAM_HEAD_TO_HEAD:
    hh = SARVAM_HEAD_TO_HEAD
    s = clone(prs, A_BARS3)
    sh = S(s)
    set_text(sh[0], "CATEGORY, HEAD TO HEAD")
    set_text(sh[1], hh["title"])
    set_text(
        sh[2],
        "Share of grievances where the category matches the one the officer recorded.",
    )
    top = max(hh["pipeline"], hh["sarvam"], hh["baseline"]) or 1.0
    bars = [
        (sh[3], sh[4], sh[5], "Our pipeline", hh["pipeline"]),
        (sh[6], sh[7], sh[8], "Sarvam", hh["sarvam"]),
        (sh[9], sh[10], sh[11], "Always guess the biggest", hh["baseline"]),
    ]
    for label_sh, bar_sh, num_sh, label, value in bars:
        set_text(label_sh, label)
        set_text(num_sh, f"{value * 100:.1f}%")
        width = max(1.35, 8.0 * (value / top))
        bar_sh.width = Inches(width)
        num_sh.width = Inches(width - 0.3)
    set_text(
        sh[15],
        f"n = {hh['n']:,} grievances from Sambalpur 2024. Both engines read the same documents.",
    )
    set_text(sh[12], hh["closing"])
    s.notes_slide.notes_text_frame.text = hh.get("notes", "")

# ══ 13 · Closing ════════════════════════════════════════════════════════════
s = clone(prs, A_CLOSING)
s.notes_slide.notes_text_frame.text = (
    "Close on the limits, not a summary.\n\n"
    "What we are not claiming: no officer minutes saved, no faster resolution, no satisfaction "
    "improvement, no accuracy figure confirmed by an officer, no redaction precision, no "
    "deduplication increment, no routing gain.\n\n"
    "Nearly every effect worth proving is blocked on a log line or a timer, not on a model."
)

# ── drop the original reference slides, keep the built ones in order ──────────
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst)[:n_original]:
    rId = sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)

prs.core_properties.title = "Janasunani 2.0 — technical briefing"
prs.core_properties.author = "Data, Policy and Innovation Centre"
prs.save(OUT)
print(f"Wrote {OUT} with {len(prs.slides._sldIdLst)} slides")

# ── export the notes, so the markdown cannot drift from the deck ──────────────
lines = [
    "# Janasunani 2.0 — technical briefing",
    "",
    "Speaker notes exported from the deck. Regenerated by `build_deck.py`; edit",
    "the notes there, not here.",
    "",
]
for i, slide in enumerate(prs.slides, 1):
    # Every archetype puts the eyebrow first and the title second.
    texts = [
        sh.text_frame.text.strip()
        for sh in slide.shapes
        if sh.has_text_frame and sh.text_frame.text.strip()
    ]
    heading = texts[1] if len(texts) > 1 else (texts[0] if texts else "")
    tf = slide.notes_slide.notes_text_frame
    body = tf.text.strip() if tf is not None else ""
    lines += [f"## Slide {i} — {heading}", "", body or "_No notes._", ""]
NOTES_MD.write_text("\n".join(lines), encoding="utf-8")
print(f"Wrote {NOTES_MD}")
